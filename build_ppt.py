"""
build_ppt.py
Generates Presentation.pptx from the content of PRESENTATION_CONTENT.md.
Uses python-pptx to build fully editable slides with:
  - Navy background on every slide
  - Flat sans-serif typography (Calibri)
  - One blue accent color (#0D6EFD)
  - Native PowerPoint shapes for all diagrams
  - Speaker notes on slides 3, 4, 5
  - No stock photos, no gradients, no shadows, no clip art

GIT AUTHOR FLAG:
  All 30 commits are authored solely by Madhusudhan Gharti.
  Aman Joshi and Siddhanta Chhetri do not appear in git history.
  The three-way role split in PRESENTATION_CONTENT.md was taken as-is,
  assumed to reflect an agreed offline division of work not tracked in git.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ─── Colour palette ───────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # primary text
LIGHT_GREY  = RGBColor(0xBB, 0xC7, 0xD6)   # secondary / body text
ACCENT      = RGBColor(0x0D, 0x6E, 0xFD)   # blue accent (shapes, rule line)
ACCENT_DARK = RGBColor(0x09, 0x47, 0xB5)   # darker shade for second header box
MID_GREY    = RGBColor(0x4A, 0x5E, 0x75)   # subtle shape fill
ARROW_COL   = RGBColor(0x5B, 0xC0, 0xDE)   # arrow / connector colour

# ─── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # truly blank — no placeholders


# ═══════════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════════

def set_bg(slide, colour=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(20), bold=False, italic=False,
                colour=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = font_size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name   = "Calibri"
    return txb


def add_rule(slide, top, left=Inches(0.6), width=Inches(12.13),
             colour=ACCENT, thickness=Pt(1.5)):
    """Draws a thin horizontal line."""
    connector = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.LINE is not directly usable; use rectangle 1px tall
        left, top, width, Pt(2)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = colour
    connector.line.fill.background()


def add_slide_number(slide, num, total=8):
    add_textbox(slide, f"{num} / {total}",
                Inches(12.0), Inches(7.1), Inches(1.1), Inches(0.3),
                font_size=Pt(11), colour=MID_GREY, align=PP_ALIGN.RIGHT)


def slide_title(slide, title_text, top=Inches(0.38)):
    """Adds a bold white slide title + a coloured underline rule."""
    add_textbox(slide, title_text,
                Inches(0.6), top, Inches(12.13), Inches(0.65),
                font_size=Pt(36), bold=True, colour=WHITE)
    add_rule(slide, top + Inches(0.68))


def add_bullet_block(slide, bullets, top, left=Inches(0.7),
                     width=Inches(12.0), font_size=Pt(22),
                     colour=LIGHT_GREY, line_spacing_pt=40):
    """
    bullets: list of (indent_level, text) tuples.
    indent_level 0 → •   level 1 → –
    """
    txb = slide.shapes.add_textbox(left, top, width, H - top - Inches(0.4))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        bullet_char = "•" if level == 0 else "–"
        indent      = "    " * level
        run = p.add_run()
        run.text = f"{indent}{bullet_char}  {text}"
        run.font.size  = font_size
        run.font.color.rgb = colour
        run.font.name  = "Calibri"
        run.font.bold  = False
        # line spacing via paragraph spacing
        p.space_before = Pt(4)
        p.space_after  = Pt(4)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_rect(slide, left, top, width, height,
             fill_colour=MID_GREY, line_colour=None, text="",
             font_size=Pt(15), text_colour=WHITE, bold=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size  = font_size
        run.font.color.rgb = text_colour
        run.font.name  = "Calibri"
        run.font.bold  = bold
    return shape


def add_label(slide, text, left, top, width=Inches(1.8), height=Inches(0.35),
              font_size=Pt(12), colour=LIGHT_GREY, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.color.rgb = colour
    run.font.name  = "Calibri"
    run.font.italic = True


def add_arrow_line(slide, x1, y1, x2, y2, colour=ARROW_COL, width_pt=1.8):
    """Draws a straight connector line between two points."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap
    # Use a thin rectangle as a pseudo-line (horizontal or vertical only).
    # For diagonal connectors we use pptx connector shapes.
    from pptx.util import Emu
    dx = x2 - x1
    dy = y2 - y1
    # Place as a very thin shape
    if abs(dx) > abs(dy):   # horizontal
        shape = slide.shapes.add_shape(1, min(x1,x2), y1 - Pt(1), abs(dx), Pt(2.5))
    else:                   # vertical
        shape = slide.shapes.add_shape(1, x1 - Pt(1), min(y1,y2), Pt(2.5), abs(dy))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s1)

# Accent bar on the left edge
add_rect(s1, Inches(0), Inches(0), Inches(0.18), H, fill_colour=ACCENT)

# Main title
add_textbox(s1,
    "Multiple-Client LAN Chat Application\nwith File Transfer",
    Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.0),
    font_size=Pt(42), bold=True, colour=WHITE)

# Subtitle rule
add_rule(s1, Inches(3.6), left=Inches(0.55), width=Inches(12.2))

# Course label
add_textbox(s1, "Network Programming",
            Inches(0.55), Inches(3.75), Inches(6), Inches(0.5),
            font_size=Pt(22), colour=ACCENT, bold=True)

# Team block
add_textbox(s1, "Madhusudhan Gharti   –   Server & Protocol",
            Inches(0.55), Inches(4.55), Inches(8), Inches(0.45),
            font_size=Pt(20), colour=LIGHT_GREY)
add_textbox(s1, "Aman Joshi              –   Client",
            Inches(0.55), Inches(5.05), Inches(8), Inches(0.45),
            font_size=Pt(20), colour=LIGHT_GREY)
add_textbox(s1, "Siddhanta Chhetri    –   File Transfer & Docs",
            Inches(0.55), Inches(5.55), Inches(8), Inches(0.45),
            font_size=Pt(20), colour=LIGHT_GREY)

add_slide_number(s1, 1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s2)
slide_title(s2, "Architecture Overview")

bullets_s2 = [
    (0, "Central server coordinates multiple TCP client connections."),
    (0, "Server broadcasts chat messages to all other clients."),
    (0, "File transfers use point-to-point routing directly to recipients."),
]
add_bullet_block(s2, bullets_s2, top=Inches(1.3), width=Inches(5.6))

# ── Diagram: Server in centre, 3 Clients around it ───────────────────────────
DX = Inches(6.5)   # left offset for diagram area

# Server box (centre)
SL = DX + Inches(2.4)
ST = Inches(2.8)
SW = Inches(2.1)
SH = Inches(0.75)
add_rect(s2, SL, ST, SW, SH, fill_colour=ACCENT,
         text="Server", font_size=Pt(17), bold=True)

# Client boxes
BW, BH = Inches(1.8), Inches(0.65)
# Client 1 – top-left
C1L, C1T = DX + Inches(0.2), Inches(1.3)
add_rect(s2, C1L, C1T, BW, BH, fill_colour=MID_GREY,
         text="Client 1", font_size=Pt(15))
# Client 2 – top-right
C2L, C2T = DX + Inches(5.0), Inches(1.3)
add_rect(s2, C2L, C2T, BW, BH, fill_colour=MID_GREY,
         text="Client 2", font_size=Pt(15))
# Client 3 – bottom-centre
C3L, C3T = DX + Inches(2.55), Inches(5.0)
add_rect(s2, C3L, C3T, BW, BH, fill_colour=MID_GREY,
         text="Client 3", font_size=Pt(15))

# Vertical lines to server (crude but editable straight shapes)
# C1 → Server (diagonal approximation: horizontal then vertical)
# Midpoints used for L-shaped connectors
MX_SRV = SL + SW / 2       # server centre-x
MY_SRV = ST + SH / 2       # server centre-y

def mid_cx(left, width): return left + width / 2
def mid_cy(top, height): return top + height / 2

# Line C1 → Server (L-shape: right then down)
add_arrow_line(s2, mid_cx(C1L, BW), mid_cy(C1T, BH),
                   mid_cx(C1L, BW), MY_SRV)
add_arrow_line(s2, mid_cx(C1L, BW), MY_SRV, MX_SRV, MY_SRV)

# Line C2 → Server (L-shape: left then down)
add_arrow_line(s2, mid_cx(C2L, BW), mid_cy(C2T, BH),
                   mid_cx(C2L, BW), MY_SRV)
add_arrow_line(s2, mid_cx(C2L, BW), MY_SRV, MX_SRV, MY_SRV)

# Line C3 → Server (vertical up)
add_arrow_line(s2, MX_SRV, ST + SH, MX_SRV, C3T)

# Labels
add_label(s2, "Broadcast (Chat)", DX + Inches(0.0), Inches(2.55),
          width=Inches(2.1), colour=ARROW_COL)
add_label(s2, "Point-to-Point (File)", DX + Inches(4.5), Inches(2.55),
          width=Inches(2.3), colour=ACCENT)

add_slide_number(s2, 2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Server & Protocol
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s3)
slide_title(s3, "Server & Protocol")

bullets_s3 = [
    (0, "Protocol uses a fixed 8-byte Header with payload length."),
    (0, "Atomic duplicate-username checks occur inside a single mutex lock."),
    (0, "Broadcasts use a copy-then-release pattern to prevent locking."),
]
add_bullet_block(s3, bullets_s3, top=Inches(1.3), width=Inches(6.0))

# ── Diagram: 8-byte header layout ────────────────────────────────────────────
HL = Inches(7.3)
HT = Inches(2.2)
HH = Inches(1.1)
HW = Inches(2.6)

add_rect(s3, HL,        HT, HW, HH, fill_colour=ACCENT,
         line_colour=WHITE,
         text="type\n(int32 · 4 bytes)", font_size=Pt(15), bold=True)
add_rect(s3, HL + HW,   HT, HW, HH, fill_colour=ACCENT_DARK,
         line_colour=WHITE,
         text="length\n(int32 · 4 bytes)", font_size=Pt(15), bold=True)

add_label(s3, "← Fixed 8-byte Header →",
          HL - Inches(0.2), HT + HH + Inches(0.05),
          width=Inches(5.5), colour=LIGHT_GREY, font_size=Pt(13))

# Payload box
add_rect(s3, HL, HT + HH + Inches(0.45), HW * 2, Inches(0.8),
         fill_colour=MID_GREY, line_colour=LIGHT_GREY,
         text="Payload  (length bytes — chat string, binary chunk, etc.)",
         font_size=Pt(13))

# Presenter notes
set_notes(s3,
    "We designed a fixed-header protocol rather than using delimiters because "
    "scanning bytes is slow and binary files naturally contain characters like "
    "newlines. To prevent race conditions, the server checks for duplicate "
    "usernames and registers new ones under a single, continuous mutex lock. "
    "For broadcasting, the server copies target sockets locally and releases "
    "the mutex before sending, ensuring a slow client's network connection "
    "cannot freeze the entire server."
)

add_slide_number(s3, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Client Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s4)
slide_title(s4, "Client Architecture")

bullets_s4 = [
    (0, "Clients use a dual-thread design (input and receive threads)."),
    (0, "Separate threads prevent blocking I/O from causing deadlocks."),
    (0, "Ensures simultaneous keyboard typing and network message receiving."),
]
add_bullet_block(s4, bullets_s4, top=Inches(1.3), width=Inches(6.2))

# ── Diagram: two parallel thread lanes ───────────────────────────────────────
DL = Inches(7.2)
DT = Inches(1.8)
TW = Inches(2.5)
TH = Inches(1.0)
GAP = Inches(0.5)

# Thread A — Main (input)
add_rect(s4, DL, DT, TW, TH, fill_colour=ACCENT,
         text="Main Thread\n(stdin / fgets)", font_size=Pt(15), bold=True)

# Thread B — Receiver
add_rect(s4, DL + TW + GAP, DT, TW, TH, fill_colour=MID_GREY,
         text="Receiver Thread\n(recv_all / socket)", font_size=Pt(15))

# Shared "Process" box below both
add_rect(s4, DL + Inches(0.6), DT + TH + Inches(0.6),
         TW * 2 - Inches(0.3), Inches(0.75),
         fill_colour=RGBColor(0x1A, 0x36, 0x5C), line_colour=ACCENT,
         text="Shared State  (atomic_int is_running, socket fd)",
         font_size=Pt(13))

# Arrows from threads down to shared box
cx_a = DL + TW / 2
cx_b = DL + TW + GAP + TW / 2
shared_top = DT + TH + Inches(0.6)

add_arrow_line(s4, cx_a, DT + TH, cx_a, shared_top)
add_arrow_line(s4, cx_b, DT + TH, cx_b, shared_top)

add_label(s4, "Concurrent Execution", DL, DT + TH + Inches(1.55),
          width=Inches(5.4), colour=LIGHT_GREY, font_size=Pt(13))

set_notes(s4,
    "If the client operated on a single thread, it would deadlock. It would be "
    "stuck waiting for user keyboard input on stdin and unable to receive "
    "network messages, or vice versa. By utilizing two POSIX threads, one "
    "handles blocking terminal input while the other continuously listens on "
    "the TCP socket, allowing seamless real-time interaction."
)

add_slide_number(s4, 4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — File Transfer
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s5)
slide_title(s5, "File Transfer")

bullets_s5 = [
    (0, "Chunked binary transfer strictly trusts Header length per chunk."),
    (0, 'Explicit binary-mode file I/O ("rb"/"wb") prevents corruption.'),
    (0, "Disconnections mid-transfer trigger partial file cleanup via remove()."),
]
add_bullet_block(s5, bullets_s5, top=Inches(1.3), width=Inches(6.2))

# ── Diagram: chunked transfer pipeline ───────────────────────────────────────
DL = Inches(7.2)
DT = Inches(1.9)
BW2 = Inches(1.8)
BH2 = Inches(0.7)
ROW_H = Inches(1.15)

steps = [
    (ACCENT,       "FILE_START\n(metadata)"),
    (MID_GREY,     "FILE_CHUNK × N\n(≤ 1024 bytes each)"),
    (RGBColor(0x19,0x7A,0x55), "FILE_END\n(0-byte signal)"),
]
for i, (col, label) in enumerate(steps):
    by = DT + i * ROW_H
    add_rect(s5, DL, by, BW2, BH2, fill_colour=col, text=label, font_size=Pt(13))
    # Arrow down to next
    if i < len(steps) - 1:
        add_arrow_line(s5,
                       DL + BW2 / 2, by + BH2,
                       DL + BW2 / 2, by + ROW_H)

# Server relay box to the right
add_rect(s5, DL + BW2 + Inches(0.55), DT + ROW_H * 0.5,
         Inches(1.5), Inches(0.7),
         fill_colour=RGBColor(0x1A,0x36,0x5C), line_colour=ACCENT,
         text="Server\n(relay, no disk)", font_size=Pt(12))

add_arrow_line(s5,
               DL + BW2,                 DT + ROW_H + BH2 / 2,
               DL + BW2 + Inches(0.55),  DT + ROW_H + BH2 / 2)

set_notes(s5,
    "Files are sent in up to 1024-byte chunks. The receiver only writes the "
    "exact payload length specified by the Header, preventing garbage padding "
    "on the final chunk. We explicitly open files in binary mode to stop "
    "OS-level newline translations from destroying compiled binaries or images. "
    "If a sender disconnects prematurely, the server synthesizes a file end "
    "signal, prompting the receiver to safely delete the corrupted partial file."
)

add_slide_number(s5, 5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Live Demo
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s6)
slide_title(s6, "Live Demo")

bullets_s6 = [
    (0, "Multi-client chat broadcasting."),
    (0, "Duplicate username rejection (interactive prompt without crashing)."),
    (0, "Direct file transfer with 100% integrity verification (cmp)."),
]
add_bullet_block(s6, bullets_s6, top=Inches(1.5), font_size=Pt(26))

add_slide_number(s6, 6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Known Limitations
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s7)
slide_title(s7, "Known Limitations")

# Two-column table layout using shapes
limits = [
    ("Max Client Cap",  "Limited to 100 concurrent connections."),
    ("Username Length", "Constrained to maximum 32 characters."),
    ("Chunk Size",      "Payloads transmitted in up to 1024 bytes."),
    ("No Encryption",   "Plaintext TCP sockets susceptible to packet sniffing."),
    ("Online Routing",  "Target users must be online for file transfers."),
]

ROW_H7  = Inches(0.75)
TBL_TOP = Inches(1.45)
TBL_L   = Inches(0.6)
COL1_W  = Inches(3.2)
COL2_W  = Inches(8.9)

for i, (label, detail) in enumerate(limits):
    row_top = TBL_TOP + i * ROW_H7
    bg_col  = RGBColor(0x13,0x24,0x3A) if i % 2 == 0 else RGBColor(0x0D,0x1B,0x2A)
    # Row background
    add_rect(s7, TBL_L, row_top, COL1_W + COL2_W, ROW_H7 - Inches(0.05),
             fill_colour=bg_col)
    # Label (accent colour, bold)
    add_textbox(s7, label,
                TBL_L + Inches(0.15), row_top + Inches(0.18),
                COL1_W - Inches(0.2), ROW_H7 - Inches(0.2),
                font_size=Pt(18), bold=True, colour=ACCENT)
    # Detail (light grey)
    add_textbox(s7, detail,
                TBL_L + COL1_W + Inches(0.15), row_top + Inches(0.18),
                COL2_W - Inches(0.2), ROW_H7 - Inches(0.2),
                font_size=Pt(18), colour=LIGHT_GREY)

add_slide_number(s7, 7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s8)

# Accent bar left edge
add_rect(s8, Inches(0), Inches(0), Inches(0.18), H, fill_colour=ACCENT)

add_textbox(s8, "Thank You!", Inches(0.55), Inches(1.6),
            Inches(12.2), Inches(1.2),
            font_size=Pt(52), bold=True, colour=WHITE)

add_textbox(s8, "Questions?", Inches(0.55), Inches(2.85),
            Inches(12.2), Inches(0.8),
            font_size=Pt(32), colour=ACCENT)

add_rule(s8, Inches(3.75), left=Inches(0.55), width=Inches(12.2))

add_textbox(s8, "Madhusudhan Gharti   –   Server & Protocol",
            Inches(0.55), Inches(4.0), Inches(8), Inches(0.45),
            font_size=Pt(20), colour=LIGHT_GREY)
add_textbox(s8, "Aman Joshi              –   Client",
            Inches(0.55), Inches(4.5), Inches(8), Inches(0.45),
            font_size=Pt(20), colour=LIGHT_GREY)
add_textbox(s8, "Siddhanta Chhetri    –   File Transfer & Docs",
            Inches(0.55), Inches(5.0), Inches(8), Inches(0.45),
            font_size=Pt(20), colour=LIGHT_GREY)

add_slide_number(s8, 8)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
OUT = "Presentation.pptx"
prs.save(OUT)
print(f"[OK] Saved: {OUT}  ({prs.slides.__len__()} slides)")
